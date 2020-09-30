from pptx import Presentation


def parse(document):
    text = []
    presentation = Presentation(document.file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n ".join(text)
